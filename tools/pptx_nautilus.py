"""NAUTILUS POTX テンプレート向け python-pptx ユーティリティ。"""

from __future__ import annotations

import re
import zipfile
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

TEMPLATE_MAIN = (
    "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
)
PRESENTATION_MAIN = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)

LAYOUT_TITLE_SLIDE = ("タイトル スライド", "title slide")
LAYOUT_TITLE_CONTENT = ("タイトルとコンテンツ", "title and content", "タイトルと内容")
LAYOUT_TITLE_PICTURE = ("タイトル付きの図", "picture with caption")
LAYOUT_CLOSING = ("最終ページ", "closing")


def potx_to_pptx_bytes(potx_path: Path) -> bytes:
    """POTX を python-pptx で開ける PPTX バイト列に変換する。"""
    out = BytesIO()
    with zipfile.ZipFile(potx_path, "r") as zin, zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                text = data.decode("utf-8")
                text = text.replace(TEMPLATE_MAIN, PRESENTATION_MAIN)
                data = text.encode("utf-8")
            zout.writestr(item, data)
    return out.getvalue()


def open_template(template: Path) -> Presentation:
    return Presentation(BytesIO(potx_to_pptx_bytes(template)))


def delete_all_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        delete_slide(prs, 0)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001
    r_id = slide_id_list[index].rId
    prs.part.drop_rel(r_id)
    del slide_id_list[index]


def find_layout(prs: Presentation, *names: str):
    lowered = [n.lower() for n in names]
    for layout in prs.slide_layouts:
        if layout.name.lower() in lowered:
            return layout
    for layout in prs.slide_layouts:
        if any(n in layout.name.lower() for n in lowered):
            return layout
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


def normalize_text(text: str) -> str:
    text = text.replace("\x0b", "\n").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def strip_md_bold(text: str) -> str:
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text)


def set_placeholder_text(slide, placeholder_types: set, text: str) -> bool:
    if not text:
        return False
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        if ph_type in placeholder_types and shape.has_text_frame:
            shape.text = text
            return True
    return False


def fill_slide(slide, title: str, body: str = "") -> None:
    if not set_placeholder_text(
        slide,
        {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE},
        title,
    ):
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title

    if body and not set_placeholder_text(
        slide,
        {
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.SUBTITLE,
        },
        body,
    ):
        for shape in slide.placeholders:
            if not shape.has_text_frame:
                continue
            ph_type = shape.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                continue
            shape.text = body
            break


def _iter_placeholders(slide, *types: PP_PLACEHOLDER):
    for shape in slide.placeholders:
        if shape.placeholder_format.type in types:
            yield shape


def find_content_placeholder(slide):
    for ph_type in (PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.PICTURE):
        for shape in _iter_placeholders(slide, ph_type):
            return shape
    return None


def find_picture_placeholder(slide):
    for shape in _iter_placeholders(slide, PP_PLACEHOLDER.PICTURE):
        return shape
    return find_content_placeholder(slide)


def find_body_placeholder(slide):
    for shape in _iter_placeholders(slide, PP_PLACEHOLDER.BODY):
        return shape
    for shape in _iter_placeholders(slide, PP_PLACEHOLDER.OBJECT):
        return shape
    return None


def placeholder_bounds(shape) -> tuple[Emu, Emu, Emu, Emu]:
    return shape.left, shape.top, shape.width, shape.height


def add_table_in_placeholder(
    slide,
    rows: list[list[str]],
    *,
    font_size: int = 10,
    header_bold: bool = True,
) -> None:
    if not rows:
        return
    cols = max(len(r) for r in rows)
    if cols == 0:
        return

    ph = find_content_placeholder(slide)
    if ph is None:
        return

    left, top, width, height = placeholder_bounds(ph)
    table_shape = slide.shapes.add_table(len(rows), cols, left, top, width, height)
    table = table_shape.table

    for r_idx, row in enumerate(rows):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            text = strip_md_bold(row[c_idx]) if c_idx < len(row) else ""
            cell.text = text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                if header_bold and r_idx == 0:
                    paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.LEFT
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_picture_in_placeholder(
    slide,
    image_path: Path,
    *,
    placeholder=None,
) -> None:
    if not image_path.is_file():
        return
    ph = placeholder or find_picture_placeholder(slide)
    if ph is None:
        slide.shapes.add_picture(str(image_path), 0, 0)
        return

    left, top, width, height = placeholder_bounds(ph)
    pic = slide.shapes.add_picture(str(image_path), left, top)
    scale = min(width / pic.width, height / pic.height, 1.0)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = left + (width - pic.width) // 2
    pic.top = top + (height - pic.height) // 2


def add_picture_at(
    slide,
    image_path: Path,
    *,
    left: Emu,
    top: Emu,
    max_width: Emu | None = None,
    max_height: Emu | None = None,
) -> None:
    if not image_path.is_file():
        return
    pic = slide.shapes.add_picture(str(image_path), left, top)
    if max_width or max_height:
        mw = max_width or pic.width
        mh = max_height or pic.height
        scale = min(mw / pic.width, mh / pic.height, 1.0)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)


def add_code_block(
    slide,
    code: str,
    *,
    font_size: int = 9,
    max_lines: int = 15,
) -> None:
    lines = code.strip().splitlines()
    if len(lines) > max_lines:
        lines = lines[:max_lines] + ["# ..."]
    text = "\n".join(lines)

    ph = find_content_placeholder(slide)
    if ph is not None and ph.has_text_frame:
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = "Consolas"
        return

    left = ph.left if ph else 0
    top = ph.top if ph else 0
    width = ph.width if ph else slide.part.presentation.slide_width // 2
    height = ph.height if ph else slide.part.presentation.slide_height // 3
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = "Consolas"


def set_speaker_notes(slide, text: str) -> None:
    if not text:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = normalize_text(strip_md_bold(text))


def fill_title_slide(slide, title: str, subtitle: str) -> None:
    fill_slide(slide, title, subtitle)


def add_slide_with_layout(prs: Presentation, *layout_names: str):
    layout = find_layout(prs, *layout_names)
    return prs.slides.add_slide(layout)


def shape_top_left(shape) -> tuple[float, float]:
    return (float(shape.top), float(shape.left))


def extract_slide_content(slide) -> tuple[str, str]:
    """スライドから (title, body) を抽出。原本スタイルは使わない。"""
    blocks: list[tuple[float, float, str, bool]] = []
    title_shape = slide.shapes.title
    title_id = title_shape.shape_id if title_shape is not None else None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        is_title = title_id is not None and shape.shape_id == title_id
        top, left = shape_top_left(shape)
        blocks.append((top, left, text, is_title))

    if not blocks:
        return "", ""

    blocks.sort(key=lambda b: (b[0], b[1]))
    title_parts: list[str] = []
    body_parts: list[str] = []

    if any(b[3] for b in blocks):
        for _, _, text, is_title in blocks:
            if is_title:
                title_parts.append(text)
            else:
                body_parts.append(text)
    else:
        first = blocks[0][2]
        if len(blocks) == 1:
            if len(first) <= 80 and "\n" not in first:
                return first, ""
            return "", first
        title_parts.append(first)
        for _, _, text, _ in blocks[1:]:
            body_parts.append(text)

    title = normalize_text("\n".join(title_parts))
    body = normalize_text("\n\n".join(body_parts))
    return title, body


def choose_layout_impl(prs: Presentation, index: int, title: str, body: str):
    if index == 0:
        return find_layout(prs, *LAYOUT_TITLE_SLIDE, "タイトルのみ")
    if not body and title:
        return find_layout(prs, "タイトルのみ", "title only", "セクション見出し")
    if body and ("|" in body or "\t" in body or re.search(r"\n.+\n.+\n", body)):
        two_col = find_layout(prs, "2 つのコンテンツ", "two content", "comparison")
        if two_col is not None:
            return two_col
    return find_layout(prs, *LAYOUT_TITLE_CONTENT)

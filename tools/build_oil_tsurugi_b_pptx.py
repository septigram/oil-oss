#!/usr/bin/env python3
"""oil-presentation-tsurugi-b.md から NAUTILUS テンプレート PPT を生成する。"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from pptx_nautilus import (
    LAYOUT_CLOSING,
    LAYOUT_TITLE_CONTENT,
    LAYOUT_TITLE_PICTURE,
    LAYOUT_TITLE_SLIDE,
    add_code_block,
    add_picture_at,
    add_picture_in_placeholder,
    add_table_in_placeholder,
    delete_all_slides,
    fill_slide,
    fill_title_slide,
    find_body_placeholder,
    find_layout,
    normalize_text,
    open_template,
    set_placeholder_text,
    set_speaker_notes,
    strip_md_bold,
)
from render_mermaid import render_mermaid

DIAGRAM_SLIDES = {4, 6, 8, 9}
DEFAULT_SOURCE = Path("docs/oil-presentation-tsurugi-b.md")
DEFAULT_TEMPLATE = Path("docs/internals/NAUTILUSテンプレート16：9_2026.potx")
DEFAULT_OUTPUT = Path("docs/presentation/ops-incident-ledger-tsurugi-b.pptx")
DEFAULT_DIAGRAM_DIR = Path("docs/presentation/work/diagrams")
PROJECT_ROOT = Path(__file__).resolve().parent.parent


@dataclass
class ParsedSlide:
    number: int
    heading: str = ""
    title_main: str = ""
    subtitle: str = ""
    sections: list[tuple[str, str]] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)
    mermaid_blocks: list[str] = field(default_factory=list)
    code_blocks: list[str] = field(default_factory=list)
    bullets: list[str] = field(default_factory=list)
    numbered_items: list[str] = field(default_factory=list)
    checklist_items: list[str] = field(default_factory=list)
    notes: str = ""
    raw_body_lines: list[str] = field(default_factory=list)


def _parse_blockquote(lines: list[str], start: int) -> tuple[str, int]:
    parts: list[str] = []
    i = start
    while i < len(lines) and lines[i].startswith(">"):
        parts.append(lines[i].lstrip(">").strip())
        i += 1
    return normalize_text("\n".join(parts)), i


def _parse_table(lines: list[str], start: int) -> tuple[list[list[str]], int]:
    rows: list[list[str]] = []
    i = start
    while i < len(lines) and "|" in lines[i]:
        line = lines[i].strip()
        if re.match(r"^\|?[\s\-:|]+\|?$", line):
            i += 1
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        rows.append(cells)
        i += 1
    return rows, i


def _parse_fenced_block(lines: list[str], start: int) -> tuple[str, str, int]:
    fence = lines[start].strip()
    lang = fence[3:].strip() or "text"
    i = start + 1
    body_lines: list[str] = []
    while i < len(lines) and not lines[i].strip().startswith("```"):
        body_lines.append(lines[i])
        i += 1
    return lang, "\n".join(body_lines).rstrip(), i + 1


def _extract_section_text(block: str, label: str) -> str:
    pattern = rf"\*\*{re.escape(label)}\*\*\s*\n+(.*?)(?=\n\*\*|\n## |\Z)"
    match = re.search(pattern, block, re.DOTALL)
    if not match:
        return ""
    text = match.group(1).strip()
    if text.startswith(">"):
        text, _ = _parse_blockquote(text.splitlines(), 0)
    return strip_md_bold(normalize_text(text))


def _extract_notes(block: str) -> str:
    notes = _extract_section_text(block, "話すポイント")
    if notes:
        return notes
    match = re.search(r"\*\*話すポイント[^*]*\*\*\s*\n+(.*?)(?=\n---|\n## |\Z)", block, re.DOTALL)
    if match:
        return strip_md_bold(normalize_text(match.group(1).strip()))
    return ""


def _extract_numbered_items(block: str) -> list[str]:
    items: list[str] = []
    for line in block.splitlines():
        m = re.match(r"^\d+\.\s+(.+)$", line.strip())
        if m:
            items.append(strip_md_bold(m.group(1)))
    return items


def _extract_checklist(block: str) -> list[str]:
    items: list[str] = []
    for line in block.splitlines():
        m = re.match(r"^-\s+\[\s*\]\s+(.+)$", line.strip())
        if m:
            items.append(strip_md_bold(m.group(1)))
    return items


def _collect_bullets(block: str) -> list[str]:
    bullets: list[str] = []
    for line in block.splitlines():
        m = re.match(r"^-\s+(.+)$", line.strip())
        if m:
            bullets.append(strip_md_bold(m.group(1)))
    return bullets


def _truncate_at_appendix(block: str) -> str:
    match = re.search(r"\n## 付録 ", block)
    if match:
        return block[: match.start()]
    match = re.search(r"\n## 変更履歴", block)
    if match:
        return block[: match.start()]
    return block


def parse_presentation_md(text: str) -> list[ParsedSlide]:
    parts = re.split(r"(?=^## スライド \d+ —)", text, flags=re.MULTILINE)
    slides: list[ParsedSlide] = []

    for part in parts:
        header = re.match(r"^## スライド (\d+) —", part)
        if not header:
            continue
        num = int(header.group(1))
        if num > 11:
            continue

        part = _truncate_at_appendix(part)
        slide = ParsedSlide(number=num)
        slide.notes = _extract_notes(part)
        slide.bullets = _collect_bullets(part)
        slide.numbered_items = _extract_numbered_items(
            _extract_section_text(part, "3 つの takeaway")
        )
        slide.checklist_items = _extract_checklist(
            _extract_section_text(part, "再現チェックリスト（Tsurugi 中心）")
            or _extract_section_text(part, "再現チェックリスト")
        )

        if num == 1:
            slide.title_main = _extract_section_text(part, "タイトル案")
            slide.subtitle = _extract_section_text(part, "サブタイトル")
            extra = _extract_section_text(part, "掲載要素")
            if extra:
                slide.sections.append(("掲載", extra))
        else:
            slide.heading = _extract_section_text(part, "見出し")

        lines = part.splitlines()
        i = 0
        current_section = ""
        while i < len(lines):
            line = lines[i]
            sec = re.match(r"^\*\*(.+?)\*\*\s*$", line.strip())
            if sec and sec.group(1) not in (
                "見出し",
                "タイトル案",
                "サブタイトル",
                "話すポイント",
                "話すポイント（30 秒）",
                "ビジュアル案",
            ):
                current_section = sec.group(1)
                i += 1
                continue

            if line.strip().startswith("|"):
                table, i = _parse_table(lines, i)
                if table:
                    slide.tables.append(table)
                continue

            if line.strip().startswith("```"):
                lang, body, i = _parse_fenced_block(lines, i)
                if lang == "mermaid":
                    slide.mermaid_blocks.append(body)
                elif lang in ("python", "text", ""):
                    slide.code_blocks.append(body)
                continue

            if line.strip().startswith(">") and "**" not in line:
                quote, i = _parse_blockquote(lines, i)
                if quote and current_section:
                    slide.sections.append((current_section, quote))
                continue

            i += 1

        slides.append(slide)

    slides.sort(key=lambda s: s.number)
    return slides


def _table_to_bullets(rows: list[list[str]], max_rows: int = 6) -> str:
    if not rows:
        return ""
    header = rows[0]
    lines: list[str] = []
    for row in rows[1 : max_rows + 1]:
        if len(row) >= 2:
            lines.append(f"• {strip_md_bold(row[0])}: {strip_md_bold(row[1])}")
        elif row:
            lines.append(f"• {strip_md_bold(row[0])}")
    if len(rows) > max_rows + 1:
        lines.append("• …")
    if not lines and header:
        lines.append(" / ".join(strip_md_bold(c) for c in header))
    return "\n".join(lines)


def _render_diagrams(slides: list[ParsedSlide], diagram_dir: Path, verbose: bool) -> dict[int, Path]:
    diagram_dir.mkdir(parents=True, exist_ok=True)
    paths: dict[int, Path] = {}
    for slide in slides:
        if not slide.mermaid_blocks:
            continue
        out = diagram_dir / f"slide-{slide.number:02d}.png"
        if verbose:
            print(f"Mermaid レンダリング: slide {slide.number} -> {out}")
        paths[slide.number] = render_mermaid(slide.mermaid_blocks[0], out)
    return paths


def _set_body_text(slide, text: str) -> None:
    if not text:
        return
    set_placeholder_text(slide, {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}, text)


def _add_table_or_body(slide, rows: list[list[str]], *, font_size: int = 9) -> None:
    if not rows:
        return
    if len(rows) <= 8 and max(len(r) for r in rows) <= 4:
        add_table_in_placeholder(slide, rows, font_size=font_size)
    else:
        _set_body_text(slide, _table_to_bullets(rows, max_rows=7))


def _build_title_slide(prs, slide, data: ParsedSlide) -> None:
    title = strip_md_bold(data.title_main) or "Tsurugi で実現する運用インシデント台帳"
    subtitle_parts = [strip_md_bold(data.subtitle)]
    for label, text in data.sections:
        subtitle_parts.append(text)
    subtitle_parts.extend([
        "発表者: （プレースホルダ）",
        "日付: 2026-07-05",
        "Tsurugi 活用例シリーズ（機能紹介は別セッション）",
    ])
    subtitle = "\n".join(p for p in subtitle_parts if p)
    fill_title_slide(slide, title, subtitle)

    img = PROJECT_ROOT / "docs/images/tsurugi-oil.png"
    if img.is_file():
        add_picture_at(
            slide,
            img,
            left=prs.slide_width - Inches(3.2),
            top=prs.slide_height - Inches(2.2),
            max_width=Inches(2.8),
            max_height=Inches(1.8),
        )


def _build_content_slide(prs, slide, data: ParsedSlide, diagram_path: Path | None) -> None:
    title = strip_md_bold(data.heading)
    fill_slide(slide, title)

    if data.number in DIAGRAM_SLIDES and diagram_path:
        add_picture_in_placeholder(slide, diagram_path)
        body_ph = find_body_placeholder(slide)
        if body_ph is not None and data.tables:
            left, top, width, height = body_ph.left, body_ph.top, body_ph.width, body_ph.height
            if len(data.tables) == 1:
                add_table_in_placeholder(slide, data.tables[0], font_size=8)
            else:
                combined: list[str] = []
                for idx, table in enumerate(data.tables[:2]):
                    combined.append(_table_to_bullets(table, max_rows=5))
                _set_body_text(slide, "\n\n".join(combined))
        return

    if data.code_blocks and data.number == 7:
        add_code_block(slide, data.code_blocks[0], font_size=8, max_lines=12)
        if len(data.tables) >= 2:
            ph = find_body_placeholder(slide)
            if ph is not None:
                rows = data.tables[1]
                add_table_in_placeholder(slide, rows, font_size=8)
        elif data.tables:
            _set_body_text(slide, _table_to_bullets(data.tables[0]))
        return

    if data.number == 5:
        if data.tables:
            _add_table_or_body(slide, data.tables[0], font_size=9)
        grafana = PROJECT_ROOT / "docs/images/tsurugi-oil-grafana.png"
        if grafana.is_file():
            add_picture_at(
                slide,
                grafana,
                left=prs.slide_width - Inches(4.5),
                top=prs.slide_height - Inches(3.0),
                max_width=Inches(4.0),
                max_height=Inches(2.5),
            )
        return

    if data.number == 11:
        parts: list[str] = []
        for i, item in enumerate(data.numbered_items, 1):
            parts.append(f"{i}. {item}")
        if data.checklist_items:
            parts.append("")
            parts.append("再現チェックリスト:")
            parts.extend(f"☐ {c}" for c in data.checklist_items)
        next_lines = [b for b in data.bullets if "oil-presentation-features" in b]
        if next_lines:
            parts.append("")
            parts.append(f"次のセッション: {next_lines[0]}")
        _set_body_text(slide, "\n".join(parts))
        return

    if data.tables:
        if len(data.tables) == 1:
            _add_table_or_body(slide, data.tables[0], font_size=9)
        else:
            parts = []
            for table in data.tables[:3]:
                parts.append(_table_to_bullets(table, max_rows=5))
            _set_body_text(slide, "\n\n".join(parts))
    elif data.bullets:
        _set_body_text(slide, "\n".join(f"• {b}" for b in data.bullets[:10]))


def build_presentation(
    source: Path,
    template: Path,
    output: Path,
    *,
    diagram_dir: Path,
    closing_slide: bool = True,
    verbose: bool = False,
) -> int:
    md_text = source.read_text(encoding="utf-8")
    slides_data = parse_presentation_md(md_text)
    if len(slides_data) < 11:
        raise RuntimeError(f"スライド数不足: {len(slides_data)} / 11")

    diagrams = _render_diagrams(slides_data, diagram_dir, verbose)

    prs = open_template(template)
    delete_all_slides(prs)

    for data in slides_data:
        if data.number == 1:
            layout = find_layout(prs, *LAYOUT_TITLE_SLIDE)
        elif data.number in DIAGRAM_SLIDES:
            layout = find_layout(prs, *LAYOUT_TITLE_PICTURE)
        else:
            layout = find_layout(prs, *LAYOUT_TITLE_CONTENT)

        slide = prs.slides.add_slide(layout)
        if data.number == 1:
            _build_title_slide(prs, slide, data)
        else:
            _build_content_slide(prs, slide, data, diagrams.get(data.number))
        set_speaker_notes(slide, data.notes)

        if verbose:
            safe_title = (data.heading or data.title_main)[:40].encode("ascii", "replace").decode()
            print(f"slide {data.number}: layout={layout.name!r} title={safe_title!r}")

    if closing_slide:
        layout = find_layout(prs, *LAYOUT_CLOSING)
        prs.slides.add_slide(layout)
        if verbose:
            print(f"slide 12: layout={layout.name!r} (closing)")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    count = len(prs.slides)
    return count


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Tsurugi 案 B プレゼン PPT を生成")
    parser.add_argument("--source", type=Path, default=DEFAULT_SOURCE)
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--diagram-dir", type=Path, default=DEFAULT_DIAGRAM_DIR)
    parser.add_argument("--closing-slide", action=argparse.BooleanOptionalAction, default=True)
    parser.add_argument("-v", "--verbose", action="store_true")
    args = parser.parse_args(argv)

    root = PROJECT_ROOT
    source = args.source if args.source.is_absolute() else root / args.source
    template = args.template if args.template.is_absolute() else root / args.template
    output = args.output if args.output.is_absolute() else root / args.output
    diagram_dir = args.diagram_dir if args.diagram_dir.is_absolute() else root / args.diagram_dir

    if not source.is_file():
        print(f"ソースが見つかりません: {source}", file=sys.stderr)
        return 1
    if not template.is_file():
        print(f"テンプレートが見つかりません: {template}", file=sys.stderr)
        return 1

    try:
        count = build_presentation(
            source,
            template,
            output,
            diagram_dir=diagram_dir,
            closing_slide=args.closing_slide,
            verbose=args.verbose,
        )
    except RuntimeError as exc:
        print(str(exc), file=sys.stderr)
        return 1

    print(f"生成完了: {count} スライド -> {output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
